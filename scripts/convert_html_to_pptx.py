"""
HTML to PowerPoint (PPTX) Converter for StellarSwap+ Pitch Deck
Converts docs/pitch-deck.html into docs/pitch-deck.pptx with rich styling, dark mode theme, and 16:9 widescreen layout.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_pitch_deck():
    # Base directory
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_file = os.path.join(base_dir, "docs", "pitch-deck.pptx")
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette matching pitch-deck.html
    BG_COLOR = RGBColor(2, 2, 5)          # #020205
    CARD_BG = RGBColor(9, 9, 11)          # #09090b
    CARD_BORDER = RGBColor(30, 30, 36)    # #1e1e24
    LIME = RGBColor(163, 230, 53)         # #a3e635
    CYAN = RGBColor(34, 211, 238)         # #22d3ee
    BLUE = RGBColor(59, 130, 246)         # #3b82f6
    AMBER = RGBColor(251, 191, 36)        # #fbbf24
    PURPLE = RGBColor(167, 139, 250)      # #a78bfa
    RED = RGBColor(248, 113, 113)         # #f87171
    WHITE = RGBColor(255, 255, 255)
    MUTED = RGBColor(148, 163, 184)       # #94a3b8

    blank_layout = prs.slide_layouts[6]

    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, badge_text, badge_color, title_prefix, title_accent, title_accent_color, slide_num):
        # Badge
        badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(4.5), Inches(0.35))
        tf_b = badge_box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        p_b = tf_b.paragraphs[0]
        p_b.text = f"●  {badge_text}"
        p_b.font.size = Pt(10.5)
        p_b.font.bold = True
        p_b.font.color.rgb = badge_color

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.78), Inches(10.5), Inches(0.65))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        r1 = p_t.add_run()
        r1.text = title_prefix
        r1.font.size = Pt(25)
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        
        if title_accent:
            r2 = p_t.add_run()
            r2.text = f" {title_accent}"
            r2.font.size = Pt(25)
            r2.font.bold = True
            r2.font.color.rgb = title_accent_color

        # Slide Number
        num_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.45), Inches(1.0), Inches(0.35))
        tf_n = num_box.text_frame
        tf_n.margin_left = tf_n.margin_top = tf_n.margin_right = tf_n.margin_bottom = 0
        p_n = tf_n.paragraphs[0]
        p_n.text = f"0{slide_num} / 09"
        p_n.alignment = PP_ALIGN.RIGHT
        p_n.font.size = Pt(10)
        p_n.font.color.rgb = MUTED

    def add_card(slide, left, top, width, height, border_color=CARD_BORDER, bg_color=CARD_BG):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.2)
        return shape

    # ==========================================================
    # SLIDE 1: Title (HTML Slide 1)
    # ==========================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)

    b_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.333), Inches(0.4))
    tf = b_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "LEVEL 5 — BLUE BELT SUBMISSION  •  STELLAR ECOSYSTEM"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN

    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(11.333), Inches(1.3))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "StellarSwap"
    r1.font.size = Pt(60)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = p.add_run()
    r2.text = "+"
    r2.font.size = Pt(60)
    r2.font.bold = True
    r2.font.color.rgb = LIME

    sub_box = s1.shapes.add_textbox(Inches(1.8), Inches(3.25), Inches(9.733), Inches(1.0))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Non-custodial DEX & Soroban Escrow Vault on Stellar — fast path payment swaps, automated trust-free escrow, and multi-wallet onboarding."
    p.font.size = Pt(15.5)
    p.font.color.rgb = MUTED

    badges = [
        ("⚡ STELLAR TESTNET", "Native Path Payments", LIME),
        ("🔒 SOROBAN RUST", "Escrow Smart Contracts", CYAN),
        ("👥 52+ USERS", "Verified Testnet Trades", AMBER),
        ("🚀 CI/CD DEPLOY", "Automated Pipelines", BLUE),
    ]
    bw = Inches(2.7)
    bh = Inches(1.1)
    bgap = Inches(0.2)
    b_left = Inches(1.0)
    b_top = Inches(4.55)

    for i, (b_title, b_sub, b_col) in enumerate(badges):
        add_card(s1, b_left + i*(bw + bgap), b_top, bw, bh, border_color=b_col)
        card_tf = s1.shapes.add_textbox(b_left + i*(bw + bgap), b_top + Inches(0.18), bw, bh - Inches(0.36)).text_frame
        card_tf.word_wrap = True
        p1 = card_tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = b_title
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = b_col
        p2 = card_tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = b_sub
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(3)

    foot_box = s1.shapes.add_textbox(Inches(1.0), Inches(6.15), Inches(11.333), Inches(0.4))
    tf = foot_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "github.com/ps910/StellarSwap-Pro   •   stellar-swap-pro.vercel.app"
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED

    sn = s1.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(2.0), Inches(0.3)).text_frame
    sn.paragraphs[0].text = "01 / 09"
    sn.paragraphs[0].font.size = Pt(10)
    sn.paragraphs[0].font.color.rgb = MUTED

    # ==========================================================
    # SLIDE 2: Problem Statement (HTML Slide 2)
    # ==========================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "PROBLEM", RED, "The DeFi UX Gap on", "Stellar", CYAN, 2)

    problems = [
        ("❌ Fragmented DEX Experience", "Users must navigate multiple platforms to swap tokens on Stellar. No single interface leverages native path payments with real-time orderbook execution. Most DEXes add unnecessary smart contract overhead.", RED),
        ("❌ No On-Chain Escrow Solution", "Stellar has no mainstream trust-free escrow service. P2P trades, freelance payments, and conditional transfers rely on off-chain trust or centralized intermediaries.", RED),
        ("❌ Wallet Fragmentation", "Each dApp only supports 1-2 wallets. Users with Freighter can't use Lobstr-only dApps. No unified multi-wallet connection layer exists for Stellar.", RED),
        ("❌ Poor Error Handling", "Most Stellar dApps show cryptic transaction errors. No contextual recovery hints, no retry logic, no distinction between wallet errors and network failures.", RED)
    ]

    cw = Inches(5.6)
    ch = Inches(2.35)
    for i, (p_title, p_desc, col) in enumerate(problems):
        row = i // 2
        col_idx = i % 2
        c_left = Inches(0.8) + col_idx * (cw + Inches(0.5))
        c_top = Inches(1.8) + row * (ch + Inches(0.35))
        
        add_card(s2, c_left, c_top, cw, ch)
        tf = s2.shapes.add_textbox(c_left + Inches(0.25), c_top + Inches(0.2), cw - Inches(0.5), ch - Inches(0.4)).text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = p_title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        p2.text = p_desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTED

    # ==========================================================
    # SLIDE 3: Solution (HTML Slide 3)
    # ==========================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "SOLUTION", LIME, "StellarSwap+ — Unified DEX &", "Escrow", LIME, 3)

    sub_box = s3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.5))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "A single production-ready dApp that combines Stellar's native path payment swap engine with a custom Soroban Rust escrow vault — all through one multi-wallet interface."
    p.font.size = Pt(13)
    p.font.color.rgb = MUTED

    solutions = [
        ("⚡ Path Payment DEX", "Native Orderbook Routing", 
         "• Execute XLM ↔ USDC swaps via Stellar's native orderbook path payments.\n• Zero smart contract vulnerability for standard swaps.\n• Best execution rate across multiple intermediate assets.\n• Sub-second route calculation with customizable slippage.", LIME),
        ("🔒 Soroban Escrow Vault", "Trustless Custody Engine",
         "• Create → Fund → Release/Refund lifecycle in Soroban Rust.\n• Time-locked asset custody with automated timeout refunds.\n• All on-chain, auditable, and decentralized on Testnet.\n• Live contract event emission synced to reactive event feed.", CYAN),
        ("👛 Multi-Wallet Kit", "Universal Onboarding",
         "• Freighter, Albedo, Lobstr, xBull, Rabet, and Demo Account.\n• One universal connection modal with web auth fallback.\n• Non-custodial signing for all operations with clear error states.\n• Horizon balance loading with dynamic reserve calculation.", AMBER)
    ]

    sw = Inches(3.64)
    sh = Inches(4.3)
    sgap = Inches(0.4)

    for i, (s_title, s_sub, s_bullets, s_col) in enumerate(solutions):
        c_left = Inches(0.8) + i * (sw + sgap)
        c_top = Inches(2.2)
        add_card(s3, c_left, c_top, sw, sh, border_color=s_col)
        
        tf = s3.shapes.add_textbox(c_left + Inches(0.25), c_top + Inches(0.25), sw - Inches(0.5), sh - Inches(0.5)).text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = s_title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = s_col
        
        p2 = tf.add_paragraph()
        p2.text = s_sub
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.space_before = Pt(4)
        
        p3 = tf.add_paragraph()
        p3.text = s_bullets
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = MUTED
        p3.space_before = Pt(12)

    # ==========================================================
    # SLIDE 4: Market Opportunity (HTML Slide 4)
    # ==========================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "MARKET", CYAN, "Market", "Opportunity", CYAN, 4)

    stats = [
        ("$500M+", "Stellar DEX Monthly Volume", LIME),
        ("8M+", "Stellar Active Accounts", CYAN),
        ("$0", "Escrow Solutions on Stellar Today", AMBER),
    ]

    for i, (val, lbl, col) in enumerate(stats):
        c_top = Inches(1.8) + i * Inches(1.6)
        add_card(s4, Inches(0.8), c_top, Inches(5.3), Inches(1.35))
        tf = s4.shapes.add_textbox(Inches(1.1), c_top + Inches(0.15), Inches(4.8), Inches(1.05)).text_frame
        p1 = tf.paragraphs[0]
        p1.text = val
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = col
        p2 = tf.add_paragraph()
        p2.text = lbl
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = MUTED

    add_card(s4, Inches(6.6), Inches(1.8), Inches(5.9), Inches(4.55))
    tf_r = s4.shapes.add_textbox(Inches(6.9), Inches(2.1), Inches(5.3), Inches(3.95)).text_frame
    tf_r.word_wrap = True

    p_h = tf_r.paragraphs[0]
    p_h.text = "Why Stellar?"
    p_h.font.size = Pt(18)
    p_h.font.bold = True
    p_h.font.color.rgb = WHITE

    reasons = [
        ("✓  3-5 second finality", "Instant transaction confirmation for DEX trades and escrow actions."),
        ("✓  ~$0.00001 per transaction", "Near-zero transaction fees make micro-swaps and escrow affordable."),
        ("✓  Native path payment for DEX aggregation", "Built-in orderbook routing finds best multi-hop conversion rates."),
        ("✓  Soroban smart contracts (Rust/WASM)", "Turing-complete, high-speed contract engine with deterministic state."),
        ("✓  Growing ecosystem & institutional adoption", "USDC & EURC anchors with multi-billion dollar real-world asset volume.")
    ]

    for r_t, r_d in reasons:
        p_t = tf_r.add_paragraph()
        p_t.space_before = Pt(9)
        r1 = p_t.add_run()
        r1.text = f"{r_t} — "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = LIME
        r2 = p_t.add_run()
        r2.text = r_d
        r2.font.size = Pt(10)
        r2.font.color.rgb = MUTED

    # ==========================================================
    # SLIDE 5: Technical Architecture (HTML Slide 5)
    # ==========================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "ARCHITECTURE", LIME, "Technical", "Architecture", LIME, 5)

    layers = [
        ("FRONTEND", [
            ("React 18 + TypeScript + Vite", "Type-safe modular application"),
            ("TailwindCSS + Custom Dark Theme", "Rich glassmorphism design system"),
            ("React.lazy() Code Splitting", "Resilient chunk loading with retry"),
            ("Sentry SDK + Web Vitals", "Real-time client monitoring & telemetry"),
            ("StellarWalletsKit Multi-Wallet", "Universal wallet adapter integration")
        ], LIME),
        ("SMART CONTRACTS", [
            ("Soroban Escrow Vault (Rust)", "Create, fund, release, & timeout refund"),
            ("Soroban AMM Swap Pool (Rust)", "Constant-product liquidity pool"),
            ("13 Unit Tests (7 + 6)", "100% test coverage for contracts"),
            ("cargo fmt + build + test CI", "Strict linting and automated checks")
        ], CYAN),
        ("CI/CD PIPELINE", [
            ("GitHub Actions — Contract CI", "WASM target build and test matrix"),
            ("GitHub Actions — Frontend CI", "Typecheck, lint, and build verification"),
            ("GitHub Actions — CD Deploy", "Automated production deployment"),
            ("Vercel Production Hosting", "Edge-optimized global CDN distribution")
        ], BLUE),
        ("INFRASTRUCTURE", [
            ("Stellar Horizon API + Soroban RPC", "Live testnet blockchain connectivity"),
            ("Content Security Policy (CSP)", "XSS & script injection protection"),
            ("Exponential Backoff Retry", "Network resilience with jitter"),
            ("ErrorBoundary + Crash Recovery", "Zero white-screens guarantee"),
            ("Environment Variables (.env)", "Production config management")
        ], AMBER)
    ]

    for i, (l_title, items, col) in enumerate(layers):
        row = i // 2
        col_idx = i % 2
        c_left = Inches(0.8) + col_idx * Inches(5.9)
        c_top = Inches(1.8) + row * Inches(2.4)
        cw = Inches(5.6)
        ch = Inches(2.2)
        
        add_card(s5, c_left, c_top, cw, ch, border_color=col)
        tf = s5.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.15), cw - Inches(0.4), ch - Inches(0.3)).text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = l_title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = col
        
        for item_t, item_d in items:
            p_item = tf.add_paragraph()
            p_item.space_before = Pt(3)
            r1 = p_item.add_run()
            r1.text = f"• {item_t}: "
            r1.font.bold = True
            r1.font.size = Pt(9.5)
            r1.font.color.rgb = WHITE
            r2 = p_item.add_run()
            r2.text = item_d
            r2.font.size = Pt(9)
            r2.font.color.rgb = MUTED

    # ==========================================================
    # SLIDE 6: Traction & Metrics (HTML Slide 6)
    # ==========================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "TRACTION", LIME, "Growth &", "Metrics", LIME, 6)

    t_stats = [
        ("52+", "Testnet Users", LIME),
        ("170+", "Transactions", CYAN),
        ("4.44", "Avg Rating (out of 5)", AMBER),
        ("99.8%", "Platform Uptime", WHITE)
    ]

    for i, (val, lbl, col) in enumerate(t_stats):
        c_left = Inches(0.8) + i * Inches(2.98)
        add_card(s6, c_left, Inches(1.6), Inches(2.8), Inches(1.15))
        tf = s6.shapes.add_textbox(c_left, Inches(1.7), Inches(2.8), Inches(0.95)).text_frame
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = val
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = col
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = lbl
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = MUTED

    add_card(s6, Inches(0.8), Inches(2.95), Inches(5.6), Inches(3.9))
    tf_q = s6.shapes.add_textbox(Inches(1.05), Inches(3.1), Inches(5.1), Inches(3.6)).text_frame
    tf_q.word_wrap = True

    p_qh = tf_q.paragraphs[0]
    p_qh.text = "USER FEEDBACK HIGHLIGHTS"
    p_qh.font.size = Pt(12)
    p_qh.font.bold = True
    p_qh.font.color.rgb = LIME

    quotes = [
        ('"Connecting via Freighter was instant, and the real-time event feed confirmed the path payment in under 2 seconds."', "User #1 (5/5)"),
        ('"The Escrow vault flow is super intuitive. Creating, funding, and releasing funds was clean."', "User #3 (5/5)"),
        ('"Loved Albedo fallback when Freighter wasn\'t installed. Zero white-screens."', "User #5 (5/5)"),
        ('"The swap was instantaneous! Extremely clean interface for swapping XLM."', "Aarav S. (5/5)")
    ]

    for q_txt, q_auth in quotes:
        p_q = tf_q.add_paragraph()
        p_q.space_before = Pt(6)
        p_q.text = q_txt
        p_q.font.size = Pt(9.5)
        p_q.font.italic = True
        p_q.font.color.rgb = WHITE
        p_a = tf_q.add_paragraph()
        p_a.text = f"— {q_auth}"
        p_a.font.size = Pt(8.5)
        p_a.font.color.rgb = CYAN

    add_card(s6, Inches(6.7), Inches(2.95), Inches(5.8), Inches(3.9))
    tf_imp = s6.shapes.add_textbox(Inches(6.95), Inches(3.1), Inches(5.3), Inches(3.6)).text_frame
    tf_imp.word_wrap = True

    p_ih = tf_imp.paragraphs[0]
    p_ih.text = "IMPROVEMENTS MADE (FROM FEEDBACK)"
    p_ih.font.size = Pt(12)
    p_ih.font.bold = True
    p_ih.font.color.rgb = CYAN

    improvements = [
        "✅ Added NPS survey + feature request picker (Commit 0088148)",
        "✅ Added Analytics Dashboard with export proof (Commit 053d027)",
        "✅ Improved onboarding with Google Form embed (Commit 4d5a948)",
        "✅ Added trust badges on landing page (Commit 671b3cd)",
        "✅ Expanded features section from 3 → 6 (Commit 671b3cd)",
        "✅ Added share/referral CTA for user growth (Commit 0088148)",
        "✅ Enhanced mobile 3-tab responsive navigation (Commit be729bd)"
    ]

    for imp_text in improvements:
        p_i = tf_imp.add_paragraph()
        p_i.space_before = Pt(5)
        p_i.text = imp_text
        p_i.font.size = Pt(9.5)
        p_i.font.color.rgb = MUTED

    # ==========================================================
    # SLIDE 7: Growth Strategy (HTML Slide 7)
    # ==========================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "GROWTH", BLUE, "Growth", "Strategy", BLUE, 7)

    pillars = [
        ("📣 Community Outreach", "Google Form onboarding, social media campaigns, Stellar Discord/Reddit engagement, and referral links for organic user acquisition.", LIME),
        ("🤝 Ecosystem Partnerships", "Integrate with Stellar anchor services, partner with Soroban dApp builders, and participate in SCF grants for ecosystem visibility.", CYAN),
        ("🔄 Product-Led Growth", "In-app feedback loops, post-transaction rating widgets, NPS surveys, and feature voting to drive retention and iteration.", AMBER),
        ("📊 Data-Driven Iteration", "Analytics dashboard, Sentry error tracking, Web Vitals monitoring, and user growth tracking to prioritize features.", BLUE),
        ("🎯 Content & Education", "Demo videos, pitch decks, blog posts, and step-by-step guides to reduce onboarding friction and build trust.", PURPLE),
        ("🏆 Incentivized Testing", "Testnet XLM faucet, Friendbot integration, and gamified milestone badges to encourage exploration and repeated usage.", WHITE)
    ]

    gw = Inches(3.64)
    gh = Inches(2.1)
    ggap_x = Inches(0.4)
    ggap_y = Inches(0.3)

    for i, (g_title, g_desc, g_col) in enumerate(pillars):
        row = i // 3
        col_idx = i % 3
        c_left = Inches(0.8) + col_idx * (gw + ggap_x)
        c_top = Inches(1.8) + row * (gh + ggap_y)
        
        add_card(s7, c_left, c_top, gw, gh, border_color=g_col)
        tf = s7.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.15), gw - Inches(0.4), gh - Inches(0.3)).text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = g_title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = g_col
        
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        p2.text = g_desc
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = MUTED

    # ==========================================================
    # SLIDE 8: Future Roadmap (HTML Slide 8)
    # ==========================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_header(s8, "ROADMAP", LIME, "Future", "Roadmap", LIME, 8)

    roadmap_phases = [
        ("Q3 2026", "Mainnet Launch", "Deploy Soroban contracts to Stellar Mainnet. Audit smart contracts. Enable real-value token swaps and escrow operations.", LIME),
        ("Q4 2026", "Advanced Features", "Multi-token escrow support, batch operations, price alerts, limit orders, and cross-asset liquidity pools.", CYAN),
        ("Q1 2027", "Mobile & SDK", "React Native mobile app, developer SDK for third-party integrations, and API for programmatic escrow creation.", BLUE),
        ("Q2 2027", "Governance & Scaling", "Community governance for fee parameters, protocol upgrades via DAO voting, and multi-chain bridge exploration.", AMBER),
        ("Q3 2027", "Enterprise Escrow", "B2B escrow-as-a-service for freelance platforms, marketplace integrations, and cross-border payment escrow.", PURPLE)
    ]

    rw = Inches(11.733)
    rh = Inches(0.88)
    rgap = Inches(0.16)

    for i, (quarter, r_title, r_desc, r_col) in enumerate(roadmap_phases):
        c_top = Inches(1.7) + i * (rh + rgap)
        add_card(s8, Inches(0.8), c_top, rw, rh, border_color=r_col)
        
        tf = s8.shapes.add_textbox(Inches(1.05), c_top + Inches(0.1), rw - Inches(0.5), rh - Inches(0.2)).text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{quarter} — "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = r_col
        
        r2 = p.add_run()
        r2.text = f"{r_title}: "
        r2.font.bold = True
        r2.font.size = Pt(11)
        r2.font.color.rgb = WHITE
        
        r3 = p.add_run()
        r3.text = r_desc
        r3.font.size = Pt(10)
        r3.font.color.rgb = MUTED

    # ==========================================================
    # SLIDE 9: Thank You / CTA (HTML Slide 9)
    # ==========================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9)

    t_box = s9.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.3))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "Let's Build the Future of\n"
    r1.font.size = Pt(38)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = p.add_run()
    r2.text = "Stellar DeFi"
    r2.font.size = Pt(38)
    r2.font.bold = True
    r2.font.color.rgb = LIME

    sub_box = s9.shapes.add_textbox(Inches(2.0), Inches(2.8), Inches(9.333), Inches(0.8))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "StellarSwap+ is ready for the next phase — Mainnet launch, advanced features, and ecosystem partnerships."
    p.font.size = Pt(14)
    p.font.color.rgb = MUTED

    cta_cards = [
        ("LIVE APP", "stellar-swap-pro.vercel.app ↗", "Production Vercel Deployment", LIME),
        ("GITHUB", "ps910/StellarSwap-Pro ↗", "Open Source Repository", CYAN),
        ("FEEDBACK FORM", "Feedback Responses ↗", "50 Verified Submissions", AMBER),
    ]

    cw = Inches(3.64)
    ch = Inches(1.8)
    cgap = Inches(0.4)

    for i, (c_tag, c_val, c_sub, c_col) in enumerate(cta_cards):
        c_left = Inches(0.8) + i * (cw + cgap)
        c_top = Inches(3.8)
        add_card(s9, c_left, c_top, cw, ch, border_color=c_col)
        
        tf = s9.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.2), cw - Inches(0.4), ch - Inches(0.4)).text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = c_tag
        p1.font.size = Pt(10.5)
        p1.font.bold = True
        p1.font.color.rgb = c_col
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        p2.text = c_val
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(4)
        p3.text = c_sub
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = MUTED

    foot_box = s9.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.333), Inches(0.6))
    tf = foot_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Built with ❤️ on Stellar • StellarSwap+ • Level 5 Blue Belt Submission"
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED

    sn = s9.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(2.0), Inches(0.3)).text_frame
    sn.paragraphs[0].text = "09 / 09"
    sn.paragraphs[0].font.size = Pt(10)
    sn.paragraphs[0].font.color.rgb = MUTED

    prs.save(output_file)
    print(f"Successfully converted docs/pitch-deck.html to {output_file}")

if __name__ == "__main__":
    create_pitch_deck()
